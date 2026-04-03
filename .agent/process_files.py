"""
File preprocessor for Office formats → Markdown + extracted images.

Converts .pptx, .docx, .xlsx, .msg/.eml, .pdf files into structured markdown
with embedded images extracted as separate .png files.

Usage:
    python process_files.py "path/to/file.pptx"
    python process_files.py "path/to/folder/"           # process all supported files in folder
    python process_files.py "path/to/folder/" --recurse  # include subfolders

Output:
    For each file, creates a subfolder with:
      - content.md    (structured markdown with text + image references)
      - img_001.png   (extracted images/charts)
      - img_002.png   ...
      - metadata.md   (source file info, extraction date)
"""
import sys
import os
import re
import json
import io
from pathlib import Path
from datetime import datetime

# Fix Windows console encoding
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

sys.path.insert(0, str(Path(__file__).parent))
from config import BRAIN_ROOT

# Supported extensions mapped to processor functions
SUPPORTED_EXTENSIONS = {
    '.pptx': 'process_pptx',
    '.docx': 'process_docx',
    '.xlsx': 'process_xlsx',
    '.msg': 'process_msg',
    '.eml': 'process_eml',
    '.pdf': 'process_pdf',
}


def ensure_output_dir(source_path: Path) -> Path:
    """Create output directory next to source file."""
    output_dir = source_path.parent / source_path.stem
    output_dir.mkdir(exist_ok=True)
    return output_dir


def save_image(image_bytes: bytes, output_dir: Path, index: int, ext: str = '.png') -> str:
    """Save image bytes to file, return relative filename."""
    filename = f"img_{index:03d}{ext}"
    filepath = output_dir / filename
    filepath.write_bytes(image_bytes)
    return filename


def write_metadata(output_dir: Path, source_path: Path, file_type: str, extra: dict = None):
    """Write metadata file for the extraction."""
    meta = {
        'source_file': source_path.name,
        'source_path': str(source_path),
        'file_type': file_type,
        'extracted_date': datetime.now().strftime('%Y-%m-%d %H:%M'),
    }
    if extra:
        meta.update(extra)

    content = "---\n"
    for k, v in meta.items():
        content += f"{k}: \"{v}\"\n"
    content += "---\n\n"
    content += f"Extracted from `{source_path.name}` on {meta['extracted_date']}\n"

    (output_dir / "metadata.md").write_text(content, encoding='utf-8')


# ============================================================================
# PowerPoint (.pptx)
# ============================================================================

def process_pptx(source_path: Path) -> Path:
    """Extract slides as markdown + images from PowerPoint."""
    from pptx import Presentation
    from pptx.util import Inches

    output_dir = ensure_output_dir(source_path)
    prs = Presentation(str(source_path))

    lines = []
    lines.append(f"# {source_path.stem}\n")
    lines.append(f"> Source: `{source_path.name}` | Slides: {len(prs.slides)}\n")

    img_index = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"\n---\n\n## Slide {slide_num}")

        # Extract slide title if present
        if slide.shapes.title and slide.shapes.title.text.strip():
            lines.append(f"### {slide.shapes.title.text.strip()}\n")

        # Process shapes
        for shape in slide.shapes:
            # Text content
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != (slide.shapes.title.text.strip() if slide.shapes.title else ''):
                        # Detect bullet level
                        level = para.level if hasattr(para, 'level') else 0
                        prefix = "  " * level + "- " if level > 0 else "- "
                        lines.append(f"{prefix}{text}")

            # Tables
            if shape.has_table:
                table = shape.table
                lines.append("")
                # Header row
                header = "| " + " | ".join(
                    cell.text.strip() for cell in table.rows[0].cells
                ) + " |"
                lines.append(header)
                lines.append("| " + " | ".join("---" for _ in table.rows[0].cells) + " |")
                # Data rows
                for row in table.rows[1:]:
                    row_text = "| " + " | ".join(
                        cell.text.strip() for cell in row.cells
                    ) + " |"
                    lines.append(row_text)
                lines.append("")

            # Images/Charts
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    img_index += 1
                    ext = '.' + image.content_type.split('/')[-1].replace('jpeg', 'jpg')
                    if ext not in ('.png', '.jpg', '.gif', '.bmp'):
                        ext = '.png'
                    filename = save_image(image.blob, output_dir, img_index, ext)
                    lines.append(f"\n![Slide {slide_num} Image]({filename})\n")
                except Exception:
                    lines.append(f"\n[Image on slide {slide_num} — extraction failed]\n")

        lines.append("")

    content = "\n".join(lines)
    (output_dir / "content.md").write_text(content, encoding='utf-8')
    write_metadata(output_dir, source_path, 'pptx', {'slide_count': str(len(prs.slides)), 'image_count': str(img_index)})

    return output_dir


# ============================================================================
# Word (.docx)
# ============================================================================

def process_docx(source_path: Path) -> Path:
    """Extract text + images from Word document."""
    import docx
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    output_dir = ensure_output_dir(source_path)
    doc = docx.Document(str(source_path))

    lines = []
    lines.append(f"# {source_path.stem}\n")
    lines.append(f"> Source: `{source_path.name}`\n")

    img_index = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue

        # Map Word heading styles to markdown
        if para.style.name.startswith('Heading'):
            try:
                level = int(para.style.name.replace('Heading ', '').replace('Heading', '1'))
            except ValueError:
                level = 1
            lines.append(f"\n{'#' * (level + 1)} {text}\n")
        elif para.style.name == 'List Bullet' or para.style.name.startswith('List'):
            lines.append(f"- {text}")
        elif para.style.name == 'List Number':
            lines.append(f"1. {text}")
        else:
            lines.append(text)

    # Extract images from document relationships
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_index += 1
                image_data = rel.target_part.blob
                content_type = rel.target_part.content_type
                ext = '.' + content_type.split('/')[-1].replace('jpeg', 'jpg')
                if ext not in ('.png', '.jpg', '.gif', '.bmp'):
                    ext = '.png'
                filename = save_image(image_data, output_dir, img_index, ext)
                lines.append(f"\n![Image {img_index}]({filename})\n")
            except Exception:
                pass

    # Extract tables
    for table_idx, table in enumerate(doc.tables, 1):
        lines.append(f"\n### Table {table_idx}\n")
        for row_idx, row in enumerate(table.rows):
            row_text = "| " + " | ".join(
                cell.text.strip().replace('\n', ' ') for cell in row.cells
            ) + " |"
            lines.append(row_text)
            if row_idx == 0:
                lines.append("| " + " | ".join("---" for _ in row.cells) + " |")
        lines.append("")

    content = "\n".join(lines)
    (output_dir / "content.md").write_text(content, encoding='utf-8')
    write_metadata(output_dir, source_path, 'docx', {'image_count': str(img_index)})

    return output_dir


# ============================================================================
# Excel (.xlsx)
# ============================================================================

def process_xlsx(source_path: Path) -> Path:
    """Extract sheets as markdown tables + chart images."""
    import openpyxl
    from openpyxl.chart import ChartBase

    output_dir = ensure_output_dir(source_path)
    wb = openpyxl.load_workbook(str(source_path), data_only=True)

    lines = []
    lines.append(f"# {source_path.stem}\n")
    lines.append(f"> Source: `{source_path.name}` | Sheets: {len(wb.sheetnames)}\n")

    img_index = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"\n## Sheet: {sheet_name}\n")

        # Get data range
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            lines.append("(empty sheet)\n")
            continue

        # Filter out completely empty rows
        rows = [r for r in rows if any(cell is not None for cell in r)]
        if not rows:
            lines.append("(empty sheet)\n")
            continue

        # Build markdown table (limit to first 100 rows for sanity)
        display_rows = rows[:100]
        max_cols = max(len(r) for r in display_rows)

        # Header
        header_cells = display_rows[0] if display_rows else []
        header = "| " + " | ".join(
            str(cell if cell is not None else '') for cell in header_cells
        ) + " |"
        lines.append(header)
        lines.append("| " + " | ".join("---" for _ in range(len(header_cells))) + " |")

        # Data rows
        for row in display_rows[1:]:
            row_text = "| " + " | ".join(
                str(cell if cell is not None else '') for cell in row
            ) + " |"
            lines.append(row_text)

        if len(rows) > 100:
            lines.append(f"\n... ({len(rows) - 100} more rows truncated)\n")

        # Extract embedded images
        for image in ws._images:
            try:
                img_index += 1
                img_data = image._data()
                filename = save_image(img_data, output_dir, img_index)
                lines.append(f"\n![{sheet_name} Chart {img_index}]({filename})\n")
            except Exception:
                pass

        lines.append("")

    content = "\n".join(lines)
    (output_dir / "content.md").write_text(content, encoding='utf-8')
    write_metadata(output_dir, source_path, 'xlsx', {
        'sheet_count': str(len(wb.sheetnames)),
        'image_count': str(img_index)
    })

    return output_dir


# ============================================================================
# Outlook Email (.msg)
# ============================================================================

def process_msg(source_path: Path) -> Path:
    """Extract email content + attachments from Outlook .msg file."""
    import extract_msg

    output_dir = ensure_output_dir(source_path)
    msg = extract_msg.Message(str(source_path))

    lines = []
    lines.append(f"# Email: {msg.subject or source_path.stem}\n")
    lines.append("## Email Metadata\n")
    lines.append(f"- **From**: {msg.sender or 'Unknown'}")
    lines.append(f"- **To**: {msg.to or 'Unknown'}")
    lines.append(f"- **CC**: {msg.cc or ''}")
    lines.append(f"- **Date**: {msg.date or 'Unknown'}")
    lines.append(f"- **Subject**: {msg.subject or 'No subject'}")
    lines.append("")

    lines.append("## Email Body\n")
    body = msg.body or ''
    lines.append(body)
    lines.append("")

    # Extract attachments
    img_index = 0
    if msg.attachments:
        lines.append("## Attachments\n")
        for att in msg.attachments:
            att_name = att.longFilename or att.shortFilename or 'unknown'
            att_ext = Path(att_name).suffix.lower()

            if att_ext in ('.png', '.jpg', '.jpeg', '.gif', '.bmp'):
                img_index += 1
                filename = save_image(att.data, output_dir, img_index, att_ext)
                lines.append(f"- ![{att_name}]({filename})")
            elif att_ext in SUPPORTED_EXTENSIONS:
                # Save attachment for separate processing
                att_path = output_dir / att_name
                att_path.write_bytes(att.data)
                lines.append(f"- [{att_name}]({att_name}) (saved — run `process_files.py` on this file for extraction)")
            else:
                att_path = output_dir / att_name
                att_path.write_bytes(att.data)
                lines.append(f"- [{att_name}]({att_name})")

    content = "\n".join(lines)
    (output_dir / "content.md").write_text(content, encoding='utf-8')
    write_metadata(output_dir, source_path, 'msg', {
        'subject': msg.subject or '',
        'sender': msg.sender or '',
        'date': str(msg.date or ''),
        'attachment_count': str(len(msg.attachments) if msg.attachments else 0)
    })

    msg.close()
    return output_dir


# ============================================================================
# Email (.eml)
# ============================================================================

def process_eml(source_path: Path) -> Path:
    """Extract email content + attachments from .eml file."""
    import email
    from email import policy

    output_dir = ensure_output_dir(source_path)

    with open(source_path, 'rb') as f:
        msg = email.message_from_binary_file(f, policy=policy.default)

    lines = []
    lines.append(f"# Email: {msg['subject'] or source_path.stem}\n")
    lines.append("## Email Metadata\n")
    lines.append(f"- **From**: {msg['from'] or 'Unknown'}")
    lines.append(f"- **To**: {msg['to'] or 'Unknown'}")
    lines.append(f"- **CC**: {msg['cc'] or ''}")
    lines.append(f"- **Date**: {msg['date'] or 'Unknown'}")
    lines.append(f"- **Subject**: {msg['subject'] or 'No subject'}")
    lines.append("")

    lines.append("## Email Body\n")
    body = msg.get_body(preferencelist=('plain', 'html'))
    if body:
        body_text = body.get_content()
        # Strip HTML if needed
        if body.get_content_type() == 'text/html':
            body_text = re.sub(r'<[^>]+>', '', body_text)
            body_text = re.sub(r'\s+', ' ', body_text)
        lines.append(body_text)
    lines.append("")

    # Extract attachments
    img_index = 0
    attachments = []
    for part in msg.iter_attachments():
        att_name = part.get_filename() or 'unknown'
        att_ext = Path(att_name).suffix.lower()
        att_data = part.get_content()
        if isinstance(att_data, str):
            att_data = att_data.encode('utf-8')

        if att_ext in ('.png', '.jpg', '.jpeg', '.gif', '.bmp'):
            img_index += 1
            filename = save_image(att_data, output_dir, img_index, att_ext)
            attachments.append(f"- ![{att_name}]({filename})")
        else:
            att_path = output_dir / att_name
            att_path.write_bytes(att_data)
            attachments.append(f"- [{att_name}]({att_name})")

    if attachments:
        lines.append("## Attachments\n")
        lines.extend(attachments)

    content = "\n".join(lines)
    (output_dir / "content.md").write_text(content, encoding='utf-8')
    write_metadata(output_dir, source_path, 'eml', {
        'subject': msg['subject'] or '',
        'sender': msg['from'] or '',
        'date': msg['date'] or '',
    })

    return output_dir


# ============================================================================
# PDF (enhanced — extract images alongside text)
# ============================================================================

def process_pdf(source_path: Path) -> Path:
    """Extract text + images from PDF. Enhances existing process_inbox.py with image extraction."""
    import PyPDF2

    output_dir = ensure_output_dir(source_path)

    lines = []
    lines.append(f"# {source_path.stem}\n")
    lines.append(f"> Source: `{source_path.name}`\n")

    img_index = 0

    with open(source_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        page_count = len(reader.pages)

        lines.append(f"> Pages: {page_count}\n")

        for page_num, page in enumerate(reader.pages, 1):
            lines.append(f"\n## Page {page_num}\n")

            # Extract text
            text = page.extract_text() or ''
            if text.strip():
                lines.append(text.strip())

            # Extract images
            if hasattr(page, 'images'):
                for img in page.images:
                    try:
                        img_index += 1
                        ext = Path(img.name).suffix or '.png'
                        if ext not in ('.png', '.jpg', '.jpeg', '.gif'):
                            ext = '.png'
                        filename = save_image(img.data, output_dir, img_index, ext)
                        lines.append(f"\n![Page {page_num} Image]({filename})\n")
                    except Exception:
                        pass

    content = "\n".join(lines)
    (output_dir / "content.md").write_text(content, encoding='utf-8')
    write_metadata(output_dir, source_path, 'pdf', {
        'page_count': str(page_count),
        'image_count': str(img_index)
    })

    return output_dir


# ============================================================================
# Main dispatcher
# ============================================================================

def process_file(file_path: Path) -> Path:
    """Process a single file, dispatching to the correct handler."""
    ext = file_path.suffix.lower()
    handler_name = SUPPORTED_EXTENSIONS.get(ext)

    if not handler_name:
        print(f"  [SKIP] Unsupported: {file_path.name} ({ext})")
        return None

    handler = globals()[handler_name]
    try:
        output_dir = handler(file_path)
        content_file = output_dir / "content.md"
        size = content_file.stat().st_size if content_file.exists() else 0
        img_count = len(list(output_dir.glob("img_*")))
        print(f"  [OK] {file_path.name} -> {output_dir.name}/ (content: {size:,} bytes, images: {img_count})")
        return output_dir
    except Exception as e:
        print(f"  [ERROR] {file_path.name}: {e}")
        return None


def process_folder(folder_path: Path, recurse: bool = False):
    """Process all supported files in a folder."""
    pattern = '**/*' if recurse else '*'
    files = sorted(folder_path.glob(pattern))
    supported = [f for f in files if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS]

    if not supported:
        print(f"No supported files found in {folder_path}")
        return

    print(f"\nProcessing {len(supported)} file(s) in {folder_path}")
    print("=" * 60)

    results = {'ok': 0, 'error': 0, 'skip': 0}
    for f in supported:
        # Skip files inside already-extracted folders
        if any(p.name == f.stem for p in f.parents):
            continue
        result = process_file(f)
        if result:
            results['ok'] += 1
        else:
            results['error'] += 1

    print(f"\nDone: {results['ok']} processed, {results['error']} errors")


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage:")
        print("  python process_files.py <file_or_folder> [--recurse]")
        print()
        print("Supported formats: " + ", ".join(SUPPORTED_EXTENSIONS.keys()))
        sys.exit(1)

    target = Path(sys.argv[1]).resolve()
    recurse = '--recurse' in sys.argv

    if target.is_file():
        process_file(target)
    elif target.is_dir():
        process_folder(target, recurse)
    else:
        print(f"Not found: {target}")
        sys.exit(1)
